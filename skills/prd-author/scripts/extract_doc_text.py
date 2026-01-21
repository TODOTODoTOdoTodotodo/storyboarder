#!/usr/bin/env python3
import argparse
import os
import sys
from typing import Iterable, List


def read_md(path: str) -> str:
    with open(path, "r", encoding="utf-8") as f:
        return f.read()


def _collapse_lines(lines: Iterable[str]) -> str:
    cleaned: List[str] = []
    last_blank = False
    for line in lines:
        if line is None:
            continue
        stripped = line.strip()
        if not stripped:
            if not last_blank:
                cleaned.append("")
                last_blank = True
            continue
        cleaned.append(stripped)
        last_blank = False
    return "\n".join(cleaned).strip()


def _format_table(rows: List[List[str]]) -> str:
    formatted = []
    for row in rows:
        cells = [cell.strip() for cell in row if cell is not None]
        if any(cells):
            formatted.append("\t".join(cells))
    return "\n".join(formatted)


def read_pdf(path: str) -> str:
    try:
        import pdfplumber  # type: ignore
    except Exception:
        pdfplumber = None

    if pdfplumber is not None:
        chunks = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                text = page.extract_text(x_tolerance=2, y_tolerance=2) or ""
                if text.strip():
                    chunks.append(text.strip())
                tables = page.extract_tables() or []
                for table in tables:
                    table_text = _format_table(table)
                    if table_text.strip():
                        chunks.append(table_text)
        return "\n\n".join(chunks)

    try:
        from PyPDF2 import PdfReader  # type: ignore
    except Exception as exc:
        raise RuntimeError("PDF extraction requires pdfplumber or PyPDF2") from exc

    reader = PdfReader(path)
    chunks = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            chunks.append(text.strip())
    return "\n\n".join(chunks)


def read_pptx(path: str) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:
        raise RuntimeError("PPTX extraction requires python-pptx") from exc

    prs = Presentation(path)
    chunks = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.has_text_frame:
                text = (shape.text or "").strip()
                if text:
                    slide_text.append(text)
            if getattr(shape, "has_table", False):
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [(cell.text or "").strip() for cell in row.cells]
                    rows.append(cells)
                table_text = _format_table(rows)
                if table_text.strip():
                    slide_text.append(table_text)
        notes = getattr(slide, "notes_slide", None)
        if notes and hasattr(notes, "notes_text_frame"):
            notes_text = (notes.notes_text_frame.text or "").strip()
            if notes_text:
                slide_text.append("[Notes]\n" + notes_text)
        if slide_text:
            chunks.append(f"[Slide {slide_index}]\n" + "\n".join(slide_text))
    return "\n\n".join(chunks)


def read_xlsx(path: str) -> str:
    try:
        import openpyxl  # type: ignore
    except Exception as exc:
        raise RuntimeError("XLSX extraction requires openpyxl") from exc

    wb = openpyxl.load_workbook(path, data_only=True)
    chunks = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            values = [("" if v is None else str(v)) for v in row]
            if any(v.strip() for v in values):
                rows.append("\t".join(values))
        if rows:
            chunks.append(f"[Sheet {sheet.title}]\n" + "\n".join(rows))
    return "\n\n".join(chunks)


def extract(path: str) -> str:
    ext = os.path.splitext(path)[1].lower()
    if ext in {".md", ".markdown"}:
        return read_md(path)
    if ext == ".pdf":
        return read_pdf(path)
    if ext in {".pptx", ".ppt"}:
        return read_pptx(path)
    if ext in {".xlsx", ".xls"}:
        return read_xlsx(path)
    raise RuntimeError(f"Unsupported file type: {ext}")


def main() -> int:
    parser = argparse.ArgumentParser(description="Extract text from documents.")
    parser.add_argument("inputs", nargs="+", help="Input file paths")
    parser.add_argument(
        "--out",
        default="-",
        help="Output file path (default: stdout)",
    )
    args = parser.parse_args()

    all_chunks: List[str] = []
    for path in args.inputs:
        if not os.path.isfile(path):
            print(f"Missing file: {path}", file=sys.stderr)
            return 1
        try:
            content = extract(path)
        except Exception as exc:
            print(f"Failed to extract {path}: {exc}", file=sys.stderr)
            return 2
        header = f"=== {os.path.basename(path)} ==="
        all_chunks.append(header)
        all_chunks.append(content)

    output = _collapse_lines("\n\n".join(all_chunks).splitlines()) + "\n"
    if args.out == "-":
        sys.stdout.write(output)
        return 0

    with open(args.out, "w", encoding="utf-8") as f:
        f.write(output)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
